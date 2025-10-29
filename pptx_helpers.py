#!/usr/bin/env python3
"""
PowerPoint-Verbesserungen für app_open_source_recovered.py
Enthält Hilfsfunktionen für bessere Bild-Extraktion und Layout-Erkennung
"""

def extract_images_from_pptx_manually(pptx_file_path):
    """
    Extrahiert Bilder direkt aus PPTX-Datei als Fallback
    Nutzt python-pptx direkt für maximale Bild-Erfassung

    Returns:
        Liste von Dictionaries mit Bild-Daten
    """
    try:
        from pptx import Presentation
        import base64

        images = []
        prs = Presentation(pptx_file_path)

        for slide_num, slide in enumerate(prs.slides, start=1):
            for shape in slide.shapes:
                # Normale Bilder (Pictures)
                if hasattr(shape, "image"):
                    try:
                        image_blob = shape.image.blob
                        image_base64 = base64.b64encode(image_blob).decode('utf-8')
                        mime_type = shape.image.content_type

                        images.append({
                            "slide": slide_num,
                            "type": "Picture",
                            "base64": image_base64,
                            "mime_type": mime_type,
                            "alt_text": shape.name or f"Bild auf Slide {slide_num}",
                            "width": shape.width if hasattr(shape, 'width') else None,
                            "height": shape.height if hasattr(shape, 'height') else None
                        })
                    except Exception as e:
                        print(f"⚠️ Bild auf Slide {slide_num} konnte nicht extrahiert werden: {e}")

                # Bilder in Gruppen
                elif shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
                    try:
                        for sub_shape in shape.shapes:
                            if hasattr(sub_shape, "image"):
                                try:
                                    image_blob = sub_shape.image.blob
                                    image_base64 = base64.b64encode(image_blob).decode('utf-8')
                                    mime_type = sub_shape.image.content_type

                                    images.append({
                                        "slide": slide_num,
                                        "type": "Picture (Grouped)",
                                        "base64": image_base64,
                                        "mime_type": mime_type,
                                        "alt_text": sub_shape.name or f"Gruppiertes Bild auf Slide {slide_num}",
                                        "width": sub_shape.width if hasattr(sub_shape, 'width') else None,
                                        "height": sub_shape.height if hasattr(sub_shape, 'height') else None
                                    })
                                except Exception:
                                    pass
                    except Exception:
                        pass

        print(f"📸 Manuelle Extraktion: {len(images)} Bilder gefunden")
        return images
    except Exception as e:
        print(f"❌ Manuelle Bild-Extraktion fehlgeschlagen: {e}")
        return []


def extract_layout_info_from_pptx(pptx_file_path):
    """
    Extrahiert Layout-Information aus PPTX für bessere Darstellung

    Returns:
        Liste von Dictionaries mit Layout-Info pro Slide
    """
    try:
        from pptx import Presentation

        layout_info = []
        prs = Presentation(pptx_file_path)

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_info = {
                "slide": slide_num,
                "layout_name": slide.slide_layout.name,
                "shape_count": len(slide.shapes),
                "has_title": False,
                "has_images": False,
                "has_tables": False,
                "title_text": None
            }

            for shape in slide.shapes:
                # Titel erkennen
                if shape.has_text_frame and shape.text.strip():
                    if hasattr(shape, 'top') and shape.top < 1000000:  # Oberer Bereich
                        slide_info["has_title"] = True
                        slide_info["title_text"] = shape.text.strip()[:100]

                # Tabellen
                if hasattr(shape, 'has_table') and shape.has_table:
                    slide_info["has_tables"] = True

                # Bilder
                if hasattr(shape, 'image'):
                    slide_info["has_images"] = True

            layout_info.append(slide_info)

        print(f"📐 Layout-Info extrahiert für {len(layout_info)} Slides")
        return layout_info
    except Exception as e:
        print(f"⚠️ Layout-Extraktion fehlgeschlagen: {e}")
        return []


def extract_headers_footers_from_pptx(pptx_file_path):
    """
    Extrahiert Header, Footer und Fußnoten aus PPTX

    PowerPoint hat folgende Elemente auf Slides:
    - Header (nur in Notes/Handouts)
    - Footer (Fußzeile auf Slides)
    - Slide Number (Seitenzahl)
    - Date (Datum)

    Returns:
        Dictionary mit Header/Footer-Informationen
    """
    try:
        from pptx import Presentation

        headers_footers = {
            "has_footer": False,
            "has_slide_numbers": False,
            "has_date": False,
            "slides": []
        }

        prs = Presentation(pptx_file_path)

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_hf = {
                "slide": slide_num,
                "footer_text": None,
                "slide_number": None,
                "date_text": None,
                "notes": None
            }

            # Slide Notes extrahieren (oft für Fußnoten genutzt)
            if slide.has_notes_slide:
                try:
                    notes_slide = slide.notes_slide
                    if notes_slide.notes_text_frame:
                        notes_text = notes_slide.notes_text_frame.text.strip()
                        if notes_text:
                            slide_hf["notes"] = notes_text
                except Exception:
                    pass

            # Footer/Header aus Slide Properties
            # PowerPoint speichert Header/Footer in HeadersFooters-Objekt
            try:
                # Zugriff auf Slide-Element (XML-Ebene)
                slide_part = slide.part

                # Footer-Platzhalter suchen
                for shape in slide.shapes:
                    # Footer ist oft in Placeholder mit bestimmtem Type
                    if hasattr(shape, 'placeholder_format'):
                        ph_type = shape.placeholder_format.type
                        # 4 = FOOTER, 12 = SLIDE_NUMBER, 13 = DATE
                        if ph_type == 4 and shape.has_text_frame:  # Footer
                            footer_text = shape.text.strip()
                            if footer_text:
                                slide_hf["footer_text"] = footer_text
                                headers_footers["has_footer"] = True
                        elif ph_type == 12 and shape.has_text_frame:  # Slide Number
                            slide_hf["slide_number"] = shape.text.strip()
                            headers_footers["has_slide_numbers"] = True
                        elif ph_type == 13 and shape.has_text_frame:  # Date
                            slide_hf["date_text"] = shape.text.strip()
                            headers_footers["has_date"] = True
            except Exception as e:
                print(f"⚠️ Footer-Extraktion für Slide {slide_num} fehlgeschlagen: {e}")

            headers_footers["slides"].append(slide_hf)

        # Zusammenfassung
        footer_count = sum(1 for s in headers_footers["slides"] if s["footer_text"])
        notes_count = sum(1 for s in headers_footers["slides"] if s["notes"])

        print(f"📝 Header/Footer extrahiert:")
        print(f"   - Footer auf {footer_count} Slides")
        print(f"   - Notes auf {notes_count} Slides")

        return headers_footers
    except Exception as e:
        print(f"❌ Header/Footer-Extraktion fehlgeschlagen: {e}")
        return {
            "has_footer": False,
            "has_slide_numbers": False,
            "has_date": False,
            "slides": []
        }


def elements_to_html_powerpoint_optimized(elements, layout_info=None, headers_footers=None):
    """
    Generiert HTML speziell für PowerPoint mit Slide-Struktur

    Args:
        elements: Liste von Unstructured-Elementen
        layout_info: Optional - Layout-Info von extract_layout_info_from_pptx
        headers_footers: Optional - Header/Footer-Info von extract_headers_footers_from_pptx

    Returns:
        HTML-String mit Slide-basierter Darstellung
    """
    html_parts = []

    # CSS für Slide-Layout
    css = '''
    <style>
        body {
            background-color: #f5f5f5 !important;
            color: #000 !important;
            font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, sans-serif;
        }
        .slide {
            border: 2px solid #333;
            margin: 20px auto;
            padding: 30px;
            background: white;
            max-width: 900px;
            page-break-after: always;
            box-shadow: 0 4px 6px rgba(0,0,0,0.1);
            border-radius: 8px;
            position: relative;
        }
        .slide-header {
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            color: white;
            padding: 12px 20px;
            margin: -30px -30px 25px -30px;
            font-weight: bold;
            font-size: 16px;
            border-radius: 6px 6px 0 0;
            display: flex;
            justify-content: space-between;
            align-items: center;
        }
        .layout-badge {
            background: rgba(255,255,255,0.2);
            padding: 4px 12px;
            border-radius: 12px;
            font-size: 12px;
        }
        .slide-title {
            color: #2c3e50;
            font-size: 32px;
            margin-bottom: 25px;
            border-bottom: 3px solid #667eea;
            padding-bottom: 15px;
            font-weight: 700;
        }
        .slide-text {
            font-size: 16px;
            line-height: 1.8;
            margin: 15px 0;
            color: #333;
        }
        .slide-list-item {
            margin: 8px 0;
            margin-left: 35px;
            font-size: 16px;
            line-height: 1.6;
        }
        .slide-image {
            text-align: center;
            margin: 30px 0;
            padding: 20px;
            background: #f9f9f9;
            border-radius: 8px;
        }
        .pptx-image {
            max-width: 100%;
            height: auto;
            border: 1px solid #ddd;
            box-shadow: 0 4px 8px rgba(0,0,0,0.1);
            border-radius: 4px;
        }
        .image-caption {
            font-size: 14px;
            color: #666;
            font-style: italic;
            margin-top: 12px;
        }
        .slide-table {
            margin: 20px 0;
            overflow-x: auto;
        }
        .slide-table table {
            width: 100%;
            border-collapse: collapse;
        }
        .slide-table th,
        .slide-table td {
            border: 1px solid #ddd;
            padding: 12px;
            text-align: left;
        }
        .slide-table th {
            background-color: #667eea;
            color: white;
            font-weight: 600;
        }
        .image-placeholder {
            background: #e3f2fd;
            border: 2px dashed #2196f3;
            padding: 20px;
            text-align: center;
            color: #1976d2;
            border-radius: 8px;
            margin: 15px 0;
        }
        .slide-footer {
            margin-top: 30px;
            padding-top: 15px;
            border-top: 1px solid #ddd;
            font-size: 12px;
            color: #666;
            display: flex;
            justify-content: space-between;
            align-items: center;
        }
        .footer-text {
            font-style: italic;
        }
        .slide-number {
            background: #667eea;
            color: white;
            padding: 4px 12px;
            border-radius: 12px;
            font-weight: bold;
        }
        .slide-notes {
            margin-top: 20px;
            padding: 15px;
            background: #fff3cd;
            border-left: 4px solid #ffc107;
            border-radius: 4px;
        }
        .notes-label {
            font-weight: bold;
            color: #856404;
            margin-bottom: 8px;
        }
        .notes-text {
            color: #856404;
            font-size: 14px;
            line-height: 1.6;
        }
    </style>
    '''

    html_parts.append(css)

    # Gruppiere Elemente nach Slide (page_number)
    slides = {}
    for element in elements:
        page = getattr(element.metadata, 'page_number', 1) if hasattr(element, 'metadata') else 1
        if page not in slides:
            slides[page] = []
        slides[page].append(element)

    # Generiere HTML pro Slide
    for slide_num in sorted(slides.keys()):
        slide_elements = slides[slide_num]

        # Layout-Info holen
        layout_name = "Standard"
        if layout_info:
            slide_layout = next((s for s in layout_info if s["slide"] == slide_num), None)
            if slide_layout:
                layout_name = slide_layout["layout_name"]

        # Header/Footer-Info holen
        slide_hf = None
        if headers_footers and headers_footers.get("slides"):
            slide_hf = next((s for s in headers_footers["slides"] if s["slide"] == slide_num), None)

        html_parts.append(f'<div class="slide" data-slide="{slide_num}">')
        html_parts.append(f'<div class="slide-header">')
        html_parts.append(f'<span>Slide {slide_num}</span>')
        html_parts.append(f'<span class="layout-badge">{layout_name}</span>')
        html_parts.append(f'</div>')
        html_parts.append('<div class="slide-content">')

        # Elemente auf Slide
        for element in slide_elements:
            element_type = type(element).__name__

            # Titel
            if element_type == "Title":
                html_parts.append(f'<h1 class="slide-title">{element.text}</h1>')

            # Text
            elif element_type in ["NarrativeText", "Text"]:
                html_parts.append(f'<p class="slide-text">{element.text}</p>')

            # Liste
            elif element_type == "ListItem":
                html_parts.append(f'<li class="slide-list-item">{element.text}</li>')

            # Tabelle
            elif element_type == "Table":
                table_html = element.metadata.text_as_html if hasattr(element.metadata, "text_as_html") else element.text
                html_parts.append(f'<div class="slide-table">{table_html}</div>')

            # Bild
            elif element_type in ["Image", "Figure", "Picture", "FigureCaption"]:
                if hasattr(element, 'metadata') and hasattr(element.metadata, 'image_base64') and element.metadata.image_base64:
                    mime_type = getattr(element.metadata, 'image_mime_type', 'image/png')
                    html_parts.append(f'''
                    <div class="slide-image">
                        <img src="data:{mime_type};base64,{element.metadata.image_base64}"
                             alt="{element.text}"
                             class="pptx-image"/>
                        <p class="image-caption">{element.text}</p>
                    </div>
                    ''')
                else:
                    html_parts.append(f'<div class="image-placeholder">🖼️ Bild: {element.text}</div>')

        html_parts.append('</div>')  # slide-content

        # Slide Notes anzeigen (falls vorhanden)
        if slide_hf and slide_hf.get("notes"):
            html_parts.append('<div class="slide-notes">')
            html_parts.append('<div class="notes-label">📝 Notizen:</div>')
            html_parts.append(f'<div class="notes-text">{slide_hf["notes"]}</div>')
            html_parts.append('</div>')

        # Footer anzeigen (falls vorhanden)
        if slide_hf and (slide_hf.get("footer_text") or slide_hf.get("slide_number") or slide_hf.get("date_text")):
            html_parts.append('<div class="slide-footer">')

            # Footer-Text
            footer_parts = []
            if slide_hf.get("footer_text"):
                footer_parts.append(f'<span class="footer-text">{slide_hf["footer_text"]}</span>')
            if slide_hf.get("date_text"):
                footer_parts.append(f'<span class="footer-text">📅 {slide_hf["date_text"]}</span>')

            if footer_parts:
                html_parts.append('<div>' + ' | '.join(footer_parts) + '</div>')
            else:
                html_parts.append('<div></div>')

            # Slide Number
            if slide_hf.get("slide_number"):
                html_parts.append(f'<span class="slide-number">{slide_hf["slide_number"]}</span>')

            html_parts.append('</div>')

        html_parts.append('</div>')  # slide

    return '\n'.join(html_parts)


# Verwendung:
# Importiere diese Funktionen in app_open_source_recovered.py:
# from pptx_helpers import (
#     extract_images_from_pptx_manually,
#     extract_layout_info_from_pptx,
#     extract_headers_footers_from_pptx,
#     elements_to_html_powerpoint_optimized
# )

